"""
Unit tests for PPT agent tool logic — no Copilot CLI, no python-pptx, no network required.

These tests validate the pure Python utility logic embedded in the PPT agent example:
  - Path validation (output dir scoping, traversal protection, extension enforcement)
  - Asset path validation (image extension, traversal, missing file)
  - Chart type validation
  - Table parameter validation (header/row length mismatch)
  - Slide index bounds validation
  - Image dimension validation

Run with:
    cd python && python -m pytest ../agents/tests/test_ppt_agent_tools.py -v
"""

import os
import os.path as osp
import tempfile
import textwrap

import pytest


# ---------------------------------------------------------------------------
# Inline replicas of validation helpers from examples/python.py
# These are extracted here so tests have zero external dependencies.
# ---------------------------------------------------------------------------

SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif"}

CHART_TYPE_MAP = {
    "bar": "BAR_CLUSTERED",
    "column": "COLUMN_CLUSTERED",
    "line": "LINE",
    "pie": "PIE",
    "area": "AREA",
}


def _resolve_pptx_path(file_path: str, output_dir: str) -> tuple[str | None, str | None]:
    """Inline replica of _resolve_pptx_path from examples/python.py."""
    if ".." in file_path:
        return None, "Path traversal is not allowed."
    if not file_path.lower().endswith(".pptx"):
        return None, "Only .pptx files are supported."
    if osp.isabs(file_path):
        return None, "Absolute paths are not allowed. Provide a path relative to PPTX_OUTPUT_DIR."

    output_dir_real = osp.realpath(output_dir)
    resolved = osp.realpath(osp.join(output_dir_real, file_path))
    if not resolved.startswith(output_dir_real + osp.sep) and resolved != output_dir_real:
        return None, "Access outside of the output directory is not allowed."

    return resolved, None


def _resolve_asset_path(asset_path: str, assets_dir: str) -> tuple[str | None, str | None]:
    """Inline replica of _resolve_asset_path from examples/python.py."""
    if ".." in asset_path:
        return None, "Path traversal is not allowed."

    ext = osp.splitext(asset_path)[1].lower()
    if ext not in SUPPORTED_IMAGE_EXTENSIONS:
        return None, f"Unsupported image format '{ext}'."

    if osp.isabs(asset_path):
        return None, "Absolute paths are not allowed. Provide a path relative to PPTX_ASSETS_DIR."

    assets_dir_real = osp.realpath(assets_dir)
    resolved = osp.realpath(osp.join(assets_dir_real, asset_path))
    if not resolved.startswith(assets_dir_real + osp.sep):
        return None, "Access outside the assets directory is not allowed."

    if not osp.exists(resolved):
        return None, f"Image file not found: {asset_path}"

    return resolved, None


def _validate_chart_type(chart_type: str) -> str | None:
    """Return error string if chart_type is not supported, else None."""
    if chart_type.lower() not in CHART_TYPE_MAP:
        return (
            f"Unsupported chart type '{chart_type}'. "
            f"Supported: {', '.join(sorted(CHART_TYPE_MAP))}"
        )
    return None


def _validate_table_params(
    headers: list[str], rows: list[list[str]]
) -> str | None:
    """Return error string if table parameters are invalid, else None."""
    if not headers:
        return "'headers' must not be empty"
    if not rows:
        return "'rows' must not be empty"
    for i, row in enumerate(rows):
        if len(row) != len(headers):
            return (
                f"Row {i} has {len(row)} values but expected {len(headers)} "
                f"(matching headers length)"
            )
    return None


def _validate_slide_index(slide_index: int, total_slides: int) -> str | None:
    """Return error string if slide_index is out of bounds, else None."""
    if slide_index < 0:
        return f"'slide_index' must be >= 0, got {slide_index}"
    if slide_index >= total_slides:
        return (
            f"slide_index {slide_index} is out of range "
            f"(presentation has {total_slides} slide(s))"
        )
    return None


def _validate_image_dimensions(
    left: float, top: float, width: float, height: float
) -> str | None:
    """Return error string if image dimensions are invalid, else None."""
    for name, val in [("left", left), ("top", top), ("width", width), ("height", height)]:
        if val < 0:
            return f"'{name}' must be non-negative, got {val}"
    if width == 0 or height == 0:
        return "'width' and 'height' must be greater than 0"
    return None


def _validate_series_lengths(
    categories: list[str], series: list[dict]
) -> str | None:
    """Return error string if any series length does not match categories, else None."""
    if not categories:
        return "'categories' must not be empty"
    if not series:
        return "'series' must not be empty"
    for s in series:
        if len(s["values"]) != len(categories):
            return (
                f"Series '{s['name']}' has {len(s['values'])} values but expected "
                f"{len(categories)} (matching categories length)"
            )
    return None


# ---------------------------------------------------------------------------
# Tests: PPTX output path validation
# ---------------------------------------------------------------------------


class TestResolvePptxPath:
    def test_relative_pptx_allowed(self, tmp_path):
        resolved, err = _resolve_pptx_path("report.pptx", str(tmp_path))
        assert err is None
        assert resolved.endswith("report.pptx")

    def test_nested_relative_pptx_allowed(self, tmp_path):
        resolved, err = _resolve_pptx_path("subdir/report.pptx", str(tmp_path))
        assert err is None

    def test_traversal_rejected(self, tmp_path):
        _, err = _resolve_pptx_path("../secret.pptx", str(tmp_path))
        assert err is not None
        assert "traversal" in err.lower()

    def test_non_pptx_extension_rejected(self, tmp_path):
        _, err = _resolve_pptx_path("report.docx", str(tmp_path))
        assert err is not None
        assert ".pptx" in err

    @pytest.mark.parametrize("ext", [".pdf", ".ppt", ".xlsx", ".txt", ".pptx.exe"])
    def test_non_pptx_extensions_all_rejected(self, tmp_path, ext):
        _, err = _resolve_pptx_path(f"file{ext}", str(tmp_path))
        assert err is not None

    def test_case_insensitive_extension(self, tmp_path):
        # .PPTX capitalised should be accepted
        resolved, err = _resolve_pptx_path("REPORT.PPTX", str(tmp_path))
        assert err is None

    def test_absolute_path_rejected(self, tmp_path):
        outside = tmp_path.parent / "outside.pptx"
        _, err = _resolve_pptx_path(str(outside), str(tmp_path))
        assert err is not None
        assert "Absolute paths are not allowed" in err

    def test_escape_via_symlink_target_blocked(self, tmp_path):
        # Regression: a symlink inside output_dir pointing outside must be blocked
        if not hasattr(os, "symlink"):
            pytest.skip("os.symlink not available on this platform")
        outside = tmp_path.parent / "outside.pptx"
        outside.write_bytes(b"")
        link_name = tmp_path / "escape_link.pptx"
        try:
            os.symlink(outside, link_name)
        except OSError:
            pytest.skip("Cannot create symlink on this platform (requires elevated privileges)")
        _, err = _resolve_pptx_path("escape_link.pptx", str(tmp_path))
        assert err is not None, "Symlink escaping output_dir must be rejected"


# ---------------------------------------------------------------------------
# Tests: Asset path validation
# ---------------------------------------------------------------------------


class TestResolveAssetPath:
    def test_valid_png_allowed(self, tmp_path):
        img = tmp_path / "logo.png"
        img.write_bytes(b"fake-png")
        _, err = _resolve_asset_path("logo.png", str(tmp_path))
        assert err is None

    @pytest.mark.parametrize("ext", [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif"])
    def test_all_supported_extensions_accepted(self, tmp_path, ext):
        img = tmp_path / f"image{ext}"
        img.write_bytes(b"data")
        _, err = _resolve_asset_path(f"image{ext}", str(tmp_path))
        assert err is None

    @pytest.mark.parametrize("ext", [".svg", ".webp", ".exe", ".py", ".pptx"])
    def test_unsupported_extensions_rejected(self, tmp_path, ext):
        _, err = _resolve_asset_path(f"image{ext}", str(tmp_path))
        assert err is not None
        assert "Unsupported image format" in err

    def test_traversal_in_asset_path_rejected(self, tmp_path):
        _, err = _resolve_asset_path("../outside.png", str(tmp_path))
        assert err is not None
        assert "traversal" in err.lower()

    def test_missing_file_returns_error(self, tmp_path):
        _, err = _resolve_asset_path("nonexistent.png", str(tmp_path))
        assert err is not None
        assert "not found" in err.lower()

    def test_nested_path_inside_assets_dir_allowed(self, tmp_path):
        sub = tmp_path / "icons"
        sub.mkdir()
        img = sub / "arrow.png"
        img.write_bytes(b"data")
        _, err = _resolve_asset_path("icons/arrow.png", str(tmp_path))
        assert err is None


# ---------------------------------------------------------------------------
# Tests: Chart type validation
# ---------------------------------------------------------------------------


class TestChartTypeValidation:
    @pytest.mark.parametrize("chart_type", ["bar", "column", "line", "pie", "area"])
    def test_supported_chart_types(self, chart_type):
        assert _validate_chart_type(chart_type) is None

    @pytest.mark.parametrize("chart_type", ["BAR", "Bar", "LINE"])
    def test_chart_types_are_case_insensitive(self, chart_type):
        # The validator lowercases the input — all valid types should pass
        assert _validate_chart_type(chart_type) is None

    @pytest.mark.parametrize("chart_type", ["scatter", "bubble", "radar", "histogram", ""])
    def test_unsupported_chart_types_rejected(self, chart_type):
        err = _validate_chart_type(chart_type)
        assert err is not None
        assert "Supported:" in err


# ---------------------------------------------------------------------------
# Tests: Table parameter validation
# ---------------------------------------------------------------------------


class TestTableParamValidation:
    def test_valid_table(self):
        headers = ["Region", "Q1", "Q2"]
        rows = [["EMEA", "100", "120"], ["AMER", "80", "90"]]
        assert _validate_table_params(headers, rows) is None

    def test_empty_headers_rejected(self):
        err = _validate_table_params([], [["a", "b"]])
        assert err is not None
        assert "headers" in err

    def test_empty_rows_rejected(self):
        err = _validate_table_params(["Col1", "Col2"], [])
        assert err is not None
        assert "rows" in err

    def test_row_length_mismatch_rejected(self):
        headers = ["A", "B", "C"]
        rows = [["1", "2"]]  # one column short
        err = _validate_table_params(headers, rows)
        assert err is not None
        assert "Row 0" in err
        assert "2" in err  # actual count
        assert "3" in err  # expected count

    def test_correct_row_length_passes(self):
        headers = ["Name", "Score"]
        rows = [["Alice", "95"], ["Bob", "87"]]
        assert _validate_table_params(headers, rows) is None

    def test_first_valid_row_second_invalid_caught(self):
        headers = ["X", "Y"]
        rows = [["1", "2"], ["3"]]  # row 1 is short
        err = _validate_table_params(headers, rows)
        assert err is not None
        assert "Row 1" in err


# ---------------------------------------------------------------------------
# Tests: Slide index bounds validation
# ---------------------------------------------------------------------------


class TestSlideIndexValidation:
    def test_valid_first_slide(self):
        assert _validate_slide_index(0, 5) is None

    def test_valid_last_slide(self):
        assert _validate_slide_index(4, 5) is None

    def test_negative_index_rejected(self):
        err = _validate_slide_index(-1, 5)
        assert err is not None
        assert ">= 0" in err

    def test_index_equal_to_count_rejected(self):
        err = _validate_slide_index(5, 5)
        assert err is not None
        assert "out of range" in err

    def test_index_beyond_count_rejected(self):
        err = _validate_slide_index(10, 3)
        assert err is not None
        assert "3 slide(s)" in err

    def test_single_slide_deck_zero_valid(self):
        assert _validate_slide_index(0, 1) is None

    def test_single_slide_deck_one_rejected(self):
        err = _validate_slide_index(1, 1)
        assert err is not None


# ---------------------------------------------------------------------------
# Tests: Image dimension validation
# ---------------------------------------------------------------------------


class TestImageDimensionValidation:
    def test_valid_dimensions(self):
        assert _validate_image_dimensions(1.0, 1.5, 8.0, 5.0) is None

    def test_zero_left_allowed(self):
        assert _validate_image_dimensions(0.0, 0.0, 4.0, 3.0) is None

    def test_negative_left_rejected(self):
        err = _validate_image_dimensions(-0.1, 1.0, 4.0, 3.0)
        assert err is not None
        assert "'left'" in err

    def test_negative_top_rejected(self):
        err = _validate_image_dimensions(1.0, -1.0, 4.0, 3.0)
        assert err is not None
        assert "'top'" in err

    def test_zero_width_rejected(self):
        err = _validate_image_dimensions(1.0, 1.0, 0.0, 3.0)
        assert err is not None
        assert "width" in err

    def test_zero_height_rejected(self):
        err = _validate_image_dimensions(1.0, 1.0, 4.0, 0.0)
        assert err is not None
        assert "height" in err  # confirms the field is mentioned in the error message

    def test_negative_width_rejected(self):
        err = _validate_image_dimensions(1.0, 1.0, -2.0, 3.0)
        assert err is not None

    def test_large_valid_dimensions(self):
        # Extra-wide slide (unusual but not invalid)
        assert _validate_image_dimensions(0.0, 0.0, 13.33, 7.5) is None


# ---------------------------------------------------------------------------
# Tests: Chart series / categories validation
# ---------------------------------------------------------------------------


class TestSeriesLengthValidation:
    def test_valid_single_series(self):
        cats = ["Q1", "Q2", "Q3"]
        series = [{"name": "Revenue", "values": [1.2, 1.5, 1.9]}]
        assert _validate_series_lengths(cats, series) is None

    def test_valid_multiple_series(self):
        cats = ["Jan", "Feb", "Mar"]
        series = [
            {"name": "2023", "values": [10, 12, 11]},
            {"name": "2024", "values": [12, 15, 14]},
        ]
        assert _validate_series_lengths(cats, series) is None

    def test_empty_categories_rejected(self):
        err = _validate_series_lengths([], [{"name": "A", "values": []}])
        assert err is not None
        assert "categories" in err

    def test_empty_series_rejected(self):
        err = _validate_series_lengths(["Q1", "Q2"], [])
        assert err is not None
        assert "series" in err

    def test_series_length_mismatch_rejected(self):
        cats = ["Q1", "Q2", "Q3"]
        series = [{"name": "Revenue", "values": [1.2, 1.5]}]  # only 2 values
        err = _validate_series_lengths(cats, series)
        assert err is not None
        assert "Revenue" in err
        assert "2" in err   # actual
        assert "3" in err   # expected

    def test_second_series_mismatch_caught(self):
        cats = ["A", "B"]
        series = [
            {"name": "S1", "values": [1, 2]},
            {"name": "S2", "values": [1]},  # wrong length
        ]
        err = _validate_series_lengths(cats, series)
        assert err is not None
        assert "S2" in err


# ---------------------------------------------------------------------------
# Integration: create and inspect a real .pptx (requires python-pptx)
# ---------------------------------------------------------------------------


class TestPptxIntegration:
    """
    These tests actually create and read .pptx files.
    They are skipped automatically when python-pptx is not installed.
    """

    @pytest.fixture(autouse=True)
    def require_pptx(self):
        pytest.importorskip("pptx", reason="python-pptx not installed — skipping integration tests")

    def test_create_blank_presentation(self, tmp_path):
        from pptx import Presentation

        prs = Presentation()
        out = str(tmp_path / "blank.pptx")
        prs.save(out)
        assert osp.exists(out)

        prs2 = Presentation(out)
        assert isinstance(prs2.slides, object)

    def test_add_and_read_title_slide(self, tmp_path):
        from pptx import Presentation

        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Test Title"
        out = str(tmp_path / "title.pptx")
        prs.save(out)

        prs2 = Presentation(out)
        assert len(prs2.slides) == 1
        assert prs2.slides[0].shapes.title.text == "Test Title"

    def test_slide_count_after_multiple_adds(self, tmp_path):
        from pptx import Presentation

        prs = Presentation()
        layout = prs.slide_layouts[1]  # Title and Content
        for i in range(5):
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Slide {i}"
        out = str(tmp_path / "multi.pptx")
        prs.save(out)

        prs2 = Presentation(out)
        assert len(prs2.slides) == 5

    def test_update_placeholder_text(self, tmp_path):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Original"
        out = str(tmp_path / "update.pptx")
        prs.save(out)

        # Round-trip: open and modify
        prs2 = Presentation(out)
        prs2.slides[0].shapes.title.text = "Updated"
        prs2.save(out)

        prs3 = Presentation(out)
        assert prs3.slides[0].shapes.title.text == "Updated"

    def test_delete_slide_reduces_count(self, tmp_path):
        from pptx import Presentation

        prs = Presentation()
        layout = prs.slide_layouts[1]
        for i in range(3):
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Slide {i}"
        out = str(tmp_path / "delete.pptx")
        prs.save(out)

        # Delete slide at index 1 via XML manipulation (same as agent tool)
        prs2 = Presentation(out)
        xml_slides = prs2.slides._sldIdLst
        xml_slides.remove(xml_slides[1])
        prs2.save(out)

        prs3 = Presentation(out)
        assert len(prs3.slides) == 2
        assert prs3.slides[0].shapes.title.text == "Slide 0"
        assert prs3.slides[1].shapes.title.text == "Slide 2"

    def test_template_round_trip(self, tmp_path):
        """Creating a presentation from a saved template preserves slide count."""
        from pptx import Presentation

        # Create template with 1 pre-existing slide
        template_prs = Presentation()
        template_prs.slides.add_slide(template_prs.slide_layouts[0])
        template_path = str(tmp_path / "template.pptx")
        template_prs.save(template_path)

        # Open template and add another slide
        prs = Presentation(template_path)
        prs.slides.add_slide(prs.slide_layouts[1])
        out = str(tmp_path / "from_template.pptx")
        prs.save(out)

        prs2 = Presentation(out)
        assert len(prs2.slides) == 2
