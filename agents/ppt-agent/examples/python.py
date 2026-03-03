"""
PPT Agent — Python example

Supports creating new presentations AND updating existing .pptx files
using python-pptx (the only free, open-source library with round-trip support).

Requirements:
    pip install python-pptx

Usage:
    # Recommended (auto-installs python-pptx):
    $env:PPTX_OUTPUT_DIR="./output"; uv run --with ".\\python" --with python-pptx agents/ppt-agent/examples/python.py

    # With pre-installed python-pptx:
    PPTX_OUTPUT_DIR=./output uv run --with ".\\python" agents/ppt-agent/examples/python.py
"""
import asyncio
import importlib.util
import os
import os.path as osp
import sys
from typing import Optional

# Check for python-pptx early so the user gets a clear message instead of a
# silent tool-level error deep inside the agent loop.
if importlib.util.find_spec("pptx") is None:
    print(
        "ERROR: python-pptx is not installed.\n"
        "  Run with:  uv run --with .\\python --with python-pptx "
        "agents/ppt-agent/examples/python.py\n"
        "  Or install: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)

from pydantic import BaseModel, Field

from copilot import CopilotClient, PermissionHandler
from copilot.generated.session_events import SessionEventType
from copilot.tools import define_tool


# ---------------------------------------------------------------------------
# Per-file asyncio locks — prevent concurrent read-modify-write races
# The SDK (and LLM) may dispatch multiple tool calls in parallel; without
# serialization, two threads can load the same file simultaneously, each add
# a slide to their own in-memory copy, and the last write wins.
# ---------------------------------------------------------------------------
_file_locks: dict[str, asyncio.Lock] = {}


def _get_file_lock(path: str) -> asyncio.Lock:
    """Return (creating if needed) the asyncio.Lock for *path*."""
    if path not in _file_locks:
        _file_locks[path] = asyncio.Lock()
    return _file_locks[path]


# ---------------------------------------------------------------------------
# Path validation helpers
# ---------------------------------------------------------------------------

SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif"}


def _resolve_pptx_path(file_path: str) -> tuple[Optional[str], Optional[str]]:
    """Resolve and validate a .pptx file path. Returns (resolved_path, error_or_None)."""
    if ".." in file_path:
        return None, "Path traversal is not allowed."
    if not file_path.lower().endswith(".pptx"):
        return None, "Only .pptx files are supported."

    if osp.isabs(file_path):
        resolved = osp.realpath(file_path)
    else:
        output_dir = os.environ.get("PPTX_OUTPUT_DIR", os.getcwd())
        output_dir_real = osp.realpath(output_dir)
        resolved = osp.realpath(osp.join(output_dir_real, file_path))
        if not resolved.startswith(output_dir_real + osp.sep) and resolved != output_dir_real:
            return None, "Access outside of the output directory is not allowed."

    return resolved, None


def _resolve_asset_path(asset_path: str) -> tuple[Optional[str], Optional[str]]:
    """Resolve and validate an image/asset file path. Returns (resolved_path, error_or_None)."""
    if ".." in asset_path:
        return None, "Path traversal is not allowed."

    ext = osp.splitext(asset_path)[1].lower()
    if ext not in SUPPORTED_IMAGE_EXTENSIONS:
        return None, f"Unsupported image format '{ext}'. Supported: {', '.join(sorted(SUPPORTED_IMAGE_EXTENSIONS))}"

    if osp.isabs(asset_path):
        resolved = osp.realpath(asset_path)
    else:
        assets_dir = os.environ.get("PPTX_ASSETS_DIR", os.getcwd())
        assets_dir_real = osp.realpath(assets_dir)
        resolved = osp.realpath(osp.join(assets_dir_real, asset_path))
        if not resolved.startswith(assets_dir_real + osp.sep):
            return None, "Access outside the assets directory is not allowed."

    if not osp.exists(resolved):
        return None, f"Image file not found: {asset_path}"

    return resolved, None


def _get_slide_title(slide) -> str:
    """Extract title text from a slide shape, or return a placeholder string."""
    try:
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    return "(no title)"


def _get_slide_content_preview(slide, max_chars: int = 80) -> str:
    """Build a short preview of non-title text shapes on a slide."""
    parts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape == slide.shapes.title:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
    preview = " · ".join(parts)
    return preview[:max_chars] + "…" if len(preview) > max_chars else preview


# ---------------------------------------------------------------------------
# Tool: create_presentation
# ---------------------------------------------------------------------------


class CreatePresentationParams(BaseModel):
    output_path: str = Field(
        description="Destination .pptx file path (relative to PPTX_OUTPUT_DIR or absolute)"
    )
    template_path: Optional[str] = Field(
        default=None,
        description="Optional path to an existing .pptx to use as slide master / theme base",
    )


@define_tool(description="Create a new PowerPoint presentation, optionally from a template")
async def create_presentation(params: CreatePresentationParams) -> dict:
    from pptx import Presentation

    resolved, err = _resolve_pptx_path(params.output_path)
    if err:
        return {"error": err}

    def _run() -> dict:
        if params.template_path:
            tmpl_resolved, tmpl_err = _resolve_pptx_path(params.template_path)
            if tmpl_err:
                return {"error": f"Template path error: {tmpl_err}"}
            if not osp.exists(tmpl_resolved):
                return {"error": f"Template file not found: {params.template_path}"}
            prs = Presentation(tmpl_resolved)
        else:
            prs = Presentation()

        parent = osp.dirname(resolved)
        if parent:
            os.makedirs(parent, exist_ok=True)
        prs.save(resolved)
        return {
            "path": resolved,
            "slide_count": len(prs.slides),
            "note": "Presentation created. Use add_*_slide tools to add content.",
        }

    try:
        return await asyncio.to_thread(_run)
    except Exception as exc:
        return {"error": str(exc)}


# ---------------------------------------------------------------------------
# Tool: open_presentation
# ---------------------------------------------------------------------------


class OpenPresentationParams(BaseModel):
    file_path: str = Field(description="Path to the existing .pptx file to open and inspect")


@define_tool(description="Open an existing PowerPoint presentation and return its slide structure")
async def open_presentation(params: OpenPresentationParams) -> dict:
    from pptx import Presentation

    resolved, err = _resolve_pptx_path(params.file_path)
    if err:
        return {"error": err}
    if not osp.exists(resolved):
        return {"error": f"File not found: {params.file_path}"}

    def _run() -> dict:
        prs = Presentation(resolved)
        slides_info = [
            {
                "index": i,
                "title": _get_slide_title(slide),
                "layout": slide.slide_layout.name,
            }
            for i, slide in enumerate(prs.slides)
        ]
        return {"path": resolved, "slide_count": len(prs.slides), "slides": slides_info}

    try:
        return await asyncio.to_thread(_run)
    except Exception as exc:
        return {"error": str(exc)}


# ---------------------------------------------------------------------------
# Tool: list_slides
# ---------------------------------------------------------------------------


class ListSlidesParams(BaseModel):
    file_path: str = Field(description="Path to the .pptx file")


@define_tool(description="List all slides in a presentation with titles and content preview")
async def list_slides(params: ListSlidesParams) -> dict:
    from pptx import Presentation

    resolved, err = _resolve_pptx_path(params.file_path)
    if err:
        return {"error": err}
    if not osp.exists(resolved):
        return {"error": f"File not found: {params.file_path}"}

    def _run() -> dict:
        prs = Presentation(resolved)
        slides_info = [
            {
                "index": i,
                "title": _get_slide_title(slide),
                "content_preview": _get_slide_content_preview(slide),
            }
            for i, slide in enumerate(prs.slides)
        ]
        return {"slide_count": len(prs.slides), "slides": slides_info}

    try:
        return await asyncio.to_thread(_run)
    except Exception as exc:
        return {"error": str(exc)}


# ---------------------------------------------------------------------------
# Tool: add_text_slide
# ---------------------------------------------------------------------------


class AddTextSlideParams(BaseModel):
    file_path: str = Field(description="Path to the .pptx file to modify")
    title: str = Field(description="Slide title")
    content: str = Field(description="Body text. Use \\n to separate bullet points")
    layout: str = Field(
        default="Title and Content",
        description="Slide layout name (default: 'Title and Content')",
    )


@define_tool(description="Add a text slide with a title and body content to a presentation")
async def add_text_slide(params: AddTextSlideParams) -> dict:
    from pptx import Presentation

    resolved, err = _resolve_pptx_path(params.file_path)
    if err:
        return {"error": err}
    if not osp.exists(resolved):
        return {"error": f"File not found: {params.file_path}"}

    def _run() -> dict:
        prs = Presentation(resolved)

        # Find layout by name (fall back to index 1 — "Title and Content")
        layout = next(
            (lay for lay in prs.slide_layouts if lay.name == params.layout),
            prs.slide_layouts[1],
        )
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = params.title

        # Populate body placeholder (index 1) with bullet lines
        body_placeholder = next(
            (ph for ph in slide.placeholders if ph.placeholder_format.idx == 1),
            None,
        )
        if body_placeholder and body_placeholder.has_text_frame:
            tf = body_placeholder.text_frame
            tf.clear()
            lines = params.content.split("\n")
            for i, line in enumerate(lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    tf.add_paragraph().text = line

        prs.save(resolved)
        return {"slide_index": len(prs.slides) - 1, "title": params.title}

    async with _get_file_lock(resolved):
        try:
            return await asyncio.to_thread(_run)
        except Exception as exc:
            return {"error": str(exc)}


# ---------------------------------------------------------------------------
# Tool: add_image_slide
# ---------------------------------------------------------------------------


class AddImageSlideParams(BaseModel):
    file_path: str = Field(description="Path to the .pptx file to modify")
    title: str = Field(description="Slide title")
    image_path: str = Field(
        description="Path to the image file (relative to PPTX_ASSETS_DIR or absolute)"
    )
    left: float = Field(default=1.0, description="Left position in inches")
    top: float = Field(default=1.5, description="Top position in inches")
    width: float = Field(default=8.0, description="Width in inches")
    height: float = Field(default=5.0, description="Height in inches")


@define_tool(description="Add a slide with a title and an image to a presentation")
async def add_image_slide(params: AddImageSlideParams) -> dict:
    from pptx import Presentation
    from pptx.util import Inches

    resolved, err = _resolve_pptx_path(params.file_path)
    if err:
        return {"error": err}
    if not osp.exists(resolved):
        return {"error": f"File not found: {params.file_path}"}

    img_resolved, img_err = _resolve_asset_path(params.image_path)
    if img_err:
        return {"error": img_err}

    # Validate dimensions
    for name, val in [("left", params.left), ("top", params.top),
                      ("width", params.width), ("height", params.height)]:
        if val < 0:
            return {"error": f"'{name}' must be non-negative, got {val}"}
    if params.width == 0 or params.height == 0:
        return {"error": "'width' and 'height' must be greater than 0"}

    def _run() -> dict:
        prs = Presentation(resolved)

        # Use blank layout for full image control
        blank_layout = next(
            (lay for lay in prs.slide_layouts if lay.name == "Blank"),
            prs.slide_layouts[6],
        )
        slide = prs.slides.add_slide(blank_layout)

        # Add title text box at top
        from pptx.util import Pt
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        tf.text = params.title
        tf.paragraphs[0].runs[0].font.size = Pt(24)
        tf.paragraphs[0].runs[0].font.bold = True

        slide.shapes.add_picture(
            img_resolved,
            Inches(params.left),
            Inches(params.top),
            Inches(params.width),
            Inches(params.height),
        )
        prs.save(resolved)
        return {"slide_index": len(prs.slides) - 1, "title": params.title}

    async with _get_file_lock(resolved):
        try:
            return await asyncio.to_thread(_run)
        except Exception as exc:
            return {"error": str(exc)}


# ---------------------------------------------------------------------------
# Tool: add_table_slide
# ---------------------------------------------------------------------------


class AddTableSlideParams(BaseModel):
    file_path: str = Field(description="Path to the .pptx file to modify")
    title: str = Field(description="Slide title")
    headers: list[str] = Field(description="Column header labels")
    rows: list[list[str]] = Field(description="Table data rows (each row must match headers length)")


@define_tool(description="Add a slide with a title and a data table to a presentation")
async def add_table_slide(params: AddTableSlideParams) -> dict:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    resolved, err = _resolve_pptx_path(params.file_path)
    if err:
        return {"error": err}
    if not osp.exists(resolved):
        return {"error": f"File not found: {params.file_path}"}
    if not params.headers:
        return {"error": "'headers' must not be empty"}
    if not params.rows:
        return {"error": "'rows' must not be empty"}
    for i, row in enumerate(params.rows):
        if len(row) != len(params.headers):
            return {
                "error": (
                    f"Row {i} has {len(row)} values but expected {len(params.headers)} "
                    f"(matching headers length)"
                )
            }

    def _run() -> dict:
        prs = Presentation(resolved)

        layout = next(
            (lay for lay in prs.slide_layouts if lay.name == "Title Only"),
            prs.slide_layouts[5],
        )
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = params.title

        cols = len(params.headers)
        total_rows = len(params.rows) + 1  # +1 for header row
        table_width = Inches(9.0)
        col_width = table_width // cols
        table = slide.shapes.add_table(
            total_rows, cols, Inches(0.5), Inches(1.5), table_width, Inches(0.4 * total_rows)
        ).table

        # Style header row
        for col_idx, header in enumerate(params.headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)  # dark blue
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Fill data rows
        for row_idx, row_data in enumerate(params.rows):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = value
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                # Alternate row shading
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xDD, 0xE8, 0xF5)

        prs.save(resolved)
        return {
            "slide_index": len(prs.slides) - 1,
            "title": params.title,
            "rows": len(params.rows),
            "columns": len(params.headers),
        }

    async with _get_file_lock(resolved):
        try:
            return await asyncio.to_thread(_run)
        except Exception as exc:
            return {"error": str(exc)}


# ---------------------------------------------------------------------------
# Tool: add_chart_slide
# ---------------------------------------------------------------------------

CHART_TYPE_MAP = {
    "bar": "BAR_CLUSTERED",
    "column": "COLUMN_CLUSTERED",
    "line": "LINE",
    "pie": "PIE",
    "area": "AREA",
}


class SeriesData(BaseModel):
    name: str = Field(description="Series name (e.g. 'Revenue')")
    values: list[float] = Field(description="Numeric data values, one per category")


class AddChartSlideParams(BaseModel):
    file_path: str = Field(description="Path to the .pptx file to modify")
    title: str = Field(description="Slide title")
    chart_type: str = Field(description="Chart type: bar, column, line, pie, or area")
    categories: list[str] = Field(description="X-axis category labels")
    series: list[SeriesData] = Field(description="One or more data series")


@define_tool(description="Add a slide with a title and a chart (bar, column, line, pie, or area)")
async def add_chart_slide(params: AddChartSlideParams) -> dict:
    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    resolved, err = _resolve_pptx_path(params.file_path)
    if err:
        return {"error": err}
    if not osp.exists(resolved):
        return {"error": f"File not found: {params.file_path}"}

    chart_key = params.chart_type.lower()
    if chart_key not in CHART_TYPE_MAP:
        return {
            "error": (
                f"Unsupported chart type '{params.chart_type}'. "
                f"Supported: {', '.join(sorted(CHART_TYPE_MAP))}"
            )
        }
    if not params.categories:
        return {"error": "'categories' must not be empty"}
    if not params.series:
        return {"error": "'series' must not be empty"}
    for s in params.series:
        if len(s.values) != len(params.categories):
            return {
                "error": (
                    f"Series '{s.name}' has {len(s.values)} values but expected "
                    f"{len(params.categories)} (matching categories length)"
                )
            }

    def _run() -> dict:
        prs = Presentation(resolved)

        layout = next(
            (lay for lay in prs.slide_layouts if lay.name == "Title Only"),
            prs.slide_layouts[5],
        )
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = params.title

        chart_data = ChartData()
        chart_data.categories = params.categories
        for s in params.series:
            chart_data.add_series(s.name, tuple(s.values))

        xl_chart_type = getattr(XL_CHART_TYPE, CHART_TYPE_MAP[chart_key])
        slide.shapes.add_chart(
            xl_chart_type,
            Inches(0.5),
            Inches(1.5),
            Inches(9.0),
            Inches(5.5),
            chart_data,
        )
        prs.save(resolved)
        return {
            "slide_index": len(prs.slides) - 1,
            "title": params.title,
            "chart_type": params.chart_type,
        }

    async with _get_file_lock(resolved):
        try:
            return await asyncio.to_thread(_run)
        except Exception as exc:
            return {"error": str(exc)}


# ---------------------------------------------------------------------------
# Tool: update_slide_text
# ---------------------------------------------------------------------------


class UpdateSlideTextParams(BaseModel):
    file_path: str = Field(description="Path to the .pptx file to modify")
    slide_index: int = Field(description="Zero-based slide index")
    placeholder_index: int = Field(
        description="Zero-based placeholder index (0 = title, 1 = body, etc.)"
    )
    new_text: str = Field(description="Replacement text for the placeholder")


@define_tool(description="Update the text of a placeholder on an existing slide")
async def update_slide_text(params: UpdateSlideTextParams) -> dict:
    from pptx import Presentation

    resolved, err = _resolve_pptx_path(params.file_path)
    if err:
        return {"error": err}
    if not osp.exists(resolved):
        return {"error": f"File not found: {params.file_path}"}
    if params.slide_index < 0:
        return {"error": f"'slide_index' must be >= 0, got {params.slide_index}"}
    if params.placeholder_index < 0:
        return {"error": f"'placeholder_index' must be >= 0, got {params.placeholder_index}"}

    def _run() -> dict:
        prs = Presentation(resolved)

        if params.slide_index >= len(prs.slides):
            return {
                "error": (
                    f"slide_index {params.slide_index} is out of range "
                    f"(presentation has {len(prs.slides)} slide(s))"
                )
            }

        slide = prs.slides[params.slide_index]
        placeholder = next(
            (ph for ph in slide.placeholders
             if ph.placeholder_format.idx == params.placeholder_index),
            None,
        )
        if placeholder is None:
            available = sorted(ph.placeholder_format.idx for ph in slide.placeholders)
            return {
                "error": (
                    f"Placeholder index {params.placeholder_index} not found on slide "
                    f"{params.slide_index}. Available indices: {available}"
                )
            }
        if not placeholder.has_text_frame:
            return {
                "error": f"Placeholder {params.placeholder_index} on slide {params.slide_index} "
                         f"does not have a text frame"
            }

        placeholder.text_frame.text = params.new_text
        prs.save(resolved)
        return {
            "updated": True,
            "slide_index": params.slide_index,
            "placeholder_index": params.placeholder_index,
        }

    async with _get_file_lock(resolved):
        try:
            return await asyncio.to_thread(_run)
        except Exception as exc:
            return {"error": str(exc)}


# ---------------------------------------------------------------------------
# Tool: delete_slide
# ---------------------------------------------------------------------------


class DeleteSlideParams(BaseModel):
    file_path: str = Field(description="Path to the .pptx file to modify")
    slide_index: int = Field(description="Zero-based index of the slide to delete")


@define_tool(description="Remove a slide from a presentation by index")
async def delete_slide(params: DeleteSlideParams) -> dict:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from lxml import etree  # lxml ships with python-pptx

    resolved, err = _resolve_pptx_path(params.file_path)
    if err:
        return {"error": err}
    if not osp.exists(resolved):
        return {"error": f"File not found: {params.file_path}"}
    if params.slide_index < 0:
        return {"error": f"'slide_index' must be >= 0, got {params.slide_index}"}

    def _run() -> dict:
        prs = Presentation(resolved)
        total = len(prs.slides)
        if params.slide_index >= total:
            return {
                "error": (
                    f"slide_index {params.slide_index} is out of range "
                    f"(presentation has {total} slide(s))"
                )
            }

        xml_slides = prs.slides._sldIdLst
        slide_elem = xml_slides[params.slide_index]
        xml_slides.remove(slide_elem)
        prs.save(resolved)
        return {
            "deleted": True,
            "slide_index": params.slide_index,
            "remaining_slides": total - 1,
        }

    async with _get_file_lock(resolved):
        try:
            return await asyncio.to_thread(_run)
        except Exception as exc:
            return {"error": str(exc)}


# ---------------------------------------------------------------------------
# Tool: save_presentation
# ---------------------------------------------------------------------------


class SavePresentationParams(BaseModel):
    file_path: str = Field(description="Path to the source .pptx file")
    output_path: str = Field(description="Destination path for the saved copy")


@define_tool(description="Save a presentation to a new path (save-as / export copy)")
async def save_presentation(params: SavePresentationParams) -> dict:
    from pptx import Presentation

    src_resolved, src_err = _resolve_pptx_path(params.file_path)
    if src_err:
        return {"error": f"Source path error: {src_err}"}
    if not osp.exists(src_resolved):
        return {"error": f"Source file not found: {params.file_path}"}

    dst_resolved, dst_err = _resolve_pptx_path(params.output_path)
    if dst_err:
        return {"error": f"Output path error: {dst_err}"}

    def _run() -> dict:
        prs = Presentation(src_resolved)
        parent = osp.dirname(dst_resolved)
        if parent:
            os.makedirs(parent, exist_ok=True)
        prs.save(dst_resolved)
        return {"saved_to": dst_resolved, "slide_count": len(prs.slides)}

    async with _get_file_lock(src_resolved):
        try:
            return await asyncio.to_thread(_run)
        except Exception as exc:
            return {"error": str(exc)}


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

_ALL_TOOLS = [
    create_presentation,
    open_presentation,
    list_slides,
    add_text_slide,
    add_image_slide,
    add_table_slide,
    add_chart_slide,
    update_slide_text,
    delete_slide,
    save_presentation,
]


def _resolve_cli_path() -> tuple[Optional[str], Optional[str]]:
    """Resolve the Copilot CLI path from environment or well-known npm locations.

    Returns (cli_path, js_script_or_None).
    When the CLI is a .js file we return (node_exe_path, js_path) so the SDK
    uses the absolute node path rather than relying on PATH lookup (which can
    fail inside uv's isolated environments on Windows).
    """
    import shutil

    def _node_exe() -> Optional[str]:
        """Return the real (non-symlink) absolute path to node.exe.

        Strategy (in order):
        1. shutil.which + os.path.realpath  (resolves nvm4w reparse points)
        2. Scan %LOCALAPPDATA%\\nvm\\ versioned dirs (nvm for Windows)
        3. Scan C:\\nvm4w\\nodejs\\ reparse point via realpath (nvm4w shim)
        4. Standard installer locations
        """
        # 1. PATH-based lookup
        node = shutil.which("node") or shutil.which("node.exe")
        if node:
            real = osp.realpath(node)
            if osp.isfile(real):
                return real

        # 2. nvm for Windows: %LOCALAPPDATA%\nvm\v<version>\node.exe
        local_appdata = os.environ.get("LOCALAPPDATA", "")
        if local_appdata:
            nvm_root = osp.join(local_appdata, "nvm")
            if osp.isdir(nvm_root):
                try:
                    versions = sorted(
                        (d for d in os.listdir(nvm_root)
                         if osp.isdir(osp.join(nvm_root, d))),
                        reverse=True,  # newest first
                    )
                    for ver in versions:
                        candidate = osp.join(nvm_root, ver, "node.exe")
                        if osp.isfile(candidate):
                            return candidate
                except OSError:
                    pass

        # 3. nvm4w: resolve the well-known reparse point directly
        nvm4w_node = r"C:\nvm4w\nodejs\node.EXE"
        if osp.exists(nvm4w_node):
            real = osp.realpath(nvm4w_node)
            if osp.isfile(real):
                return real

        # 4. Standard installer paths
        for candidate in [
            r"C:\Program Files\nodejs\node.exe",
            r"C:\Program Files (x86)\nodejs\node.exe",
        ]:
            if osp.isfile(candidate):
                return osp.realpath(candidate)

        return None

    def _use_js(js_path: str) -> tuple[Optional[str], Optional[str]]:
        """Return (node_exe_realpath, js_path) so the SDK never does a PATH lookup for 'node'."""
        node = _node_exe()
        if node:
            return node, js_path
        # Node not found — signal failure; don't silently fall back to "node"
        # (the SDK's PATH-based "node" lookup is what causes WinError 5)
        return None, None

    # 1. Explicit env var (accepts a .js file, a .bat, or a native binary)
    env_path = os.environ.get("COPILOT_CLI_PATH", "").strip().strip("{}")
    if env_path and osp.exists(env_path):
        if env_path.endswith(".js"):
            return _use_js(env_path)
        return env_path, None

    # 2. npm global install on Windows: %APPDATA%\npm\node_modules\@github\copilot\index.js
    appdata = os.environ.get("APPDATA", "")
    if appdata:
        js_path = osp.join(appdata, "npm", "node_modules", "@github", "copilot", "index.js")
        if osp.exists(js_path):
            return _use_js(js_path)

    # 3. npm global install on Unix/macOS: $(npm root -g)/@github/copilot/index.js
    npm_cmd = shutil.which("npm")
    if npm_cmd:
        try:
            import subprocess as _sp
            npm_root = _sp.check_output([npm_cmd, "root", "-g"], text=True).strip()
            js_path = osp.join(npm_root, "@github", "copilot", "index.js")
            if osp.exists(js_path):
                return _use_js(js_path)
        except Exception:
            pass

    return None, None


async def main() -> None:
    cli_path, js_script = _resolve_cli_path()
    if cli_path is None:
        print(
            "ERROR: Copilot CLI not found.\n"
            "  • Install it with:  npm install -g @github/copilot\n"
            "  • Or set COPILOT_CLI_PATH to the path of the CLI binary / index.js",
            file=sys.stderr,
        )
        sys.exit(1)

    client_opts: dict = {"cli_path": cli_path}
    if js_script:
        # cli_path is node.exe; prepend js_script so the SDK builds the full command:
        # [node_exe, js_script, --headless, ...]
        client_opts["cli_args"] = [js_script]

    client = CopilotClient(client_opts)
    await client.start()

    session = await client.create_session({
        "model": "gpt-4.1",
        "streaming": True,
        "on_permission_request": PermissionHandler.approve_all,
        "system_message": {
            "content": (
                "You are an expert PowerPoint presentation author. "
                "Use your tools to create and edit .pptx files from natural language instructions. "
                "Before modifying an existing file, call open_presentation or list_slides first. "
                "Always confirm the output path before creating or overwriting a file. "
                "After completing a set of changes, call save_presentation to persist the result. "
                "Report each action taken and end with the final file path and slide count."
            )
        },
        "tools": _ALL_TOOLS,
    })

    def handle_event(event) -> None:
        if event.type == SessionEventType.ASSISTANT_MESSAGE_DELTA:
            sys.stdout.write(event.data.delta_content)
            sys.stdout.flush()
        if event.type == SessionEventType.SESSION_IDLE:
            sys.stdout.write("\n\n")
            sys.stdout.flush()

    session.on(handle_event)

    output_dir = os.environ.get("PPTX_OUTPUT_DIR", os.getcwd())
    os.makedirs(output_dir, exist_ok=True)  # ensure the directory exists
    print("PowerPoint Agent (type 'exit' to quit)\n")
    print(f"   Output directory: {osp.abspath(output_dir)}")
    print("   Try: 'Create a 3-slide deck named demo.pptx with a title slide, agenda, and summary'\n")

    while True:
        try:
            user_input = input("You: ").strip()
        except (EOFError, KeyboardInterrupt):
            break
        if user_input.lower() == "exit":
            break
        if not user_input:
            continue
        sys.stdout.write("Agent: ")
        sys.stdout.flush()
        await session.send_and_wait({"prompt": user_input})

    await client.stop()


if __name__ == "__main__":
    asyncio.run(main())
